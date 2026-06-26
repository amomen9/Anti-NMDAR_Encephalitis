"""
Build 'Anti-NMDAR Encephalitis Modeling.pptx'.

The deck follows the section / subsection structure of BM_Report.tex and is
filled with presentation-style text drawn from the whole project (report,
README, contents_manifest.json, the base scope deck, and the Petri-net code).

Every figure is a dashed PLACEHOLDER box, not an embedded image. Each box shows
either the absolute repository path of the matching figure, or a web URL of a
suitable image, so the figures can be replaced manually afterwards.

Run:  python presentation/build_presentation.py
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------- #
# Paths
# --------------------------------------------------------------------------- #
HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
OUT = os.path.join(HERE, "Anti-NMDAR Encephalitis Modeling.pptx")


def repo(*parts):
    """Absolute Windows-style path to a repository file (shown on placeholders)."""
    return os.path.join(ROOT, *parts)


# --------------------------------------------------------------------------- #
# Figure options
# --------------------------------------------------------------------------- #
# Embed animated GIFs (PowerPoint plays them during the slide show; newer
# versions also loop them in Normal view) wherever a figure has a matching GIF.
# Set False to use the static PNGs everywhere.
USE_ANIMATED_GIFS = True

# Downscale heavy figures (big GIFs and large static images) into a cache folder
# before embedding, to keep the .pptx small. The originals under output/ are
# never modified. Set False to embed every figure at full resolution.
DOWNSCALE_BIG_IMAGES = True
GIF_MAX_WIDTH = 1280      # px: animated GIFs wider than this are shrunk
IMG_MAX_WIDTH = 1600      # px: static images wider than this are shrunk
ASSETS = os.path.join(HERE, "_assets")   # cache dir for downscaled copies


def animated(static_path, gif_path):
    """Prefer the animated GIF when it exists and animation is enabled."""
    if USE_ANIMATED_GIFS and os.path.exists(gif_path):
        return gif_path
    return static_path


# Repository figures. Those with an animated counterpart resolve to the .gif.
FIG_BASELINE_PN = animated(
    repo("output", "NMDAR", "images", "NMDAR_baseline_pn.png"),
    repo("output", "NMDAR", "animations", "NMDAR_baseline_animation.gif"))
FIG_PERTURB_PN = animated(
    repo("output", "NMDAR", "images", "NMDAR_perturbation_pn.png"),
    repo("output", "NMDAR", "animations", "NMDAR_perturbation_animation.gif"))
FIG_NMDAR_SIDE = animated(
    repo("output", "NMDAR", "images", "NMDAR_PN.png"),
    repo("output", "NMDAR", "animations", "NMDAR_PN.gif"))
FIG_CA_INFLUX = animated(
    repo("output", "NMDAR", "images", "ca_influx_comparison.png"),
    repo("output", "NMDAR", "animations", "ca_influx_comparison.gif"))
FIG_NEURO_PN = animated(
    repo("output", "Neurodegeneration", "images", "neurodegeneration_pn.png"),
    repo("output", "Neurodegeneration", "animations", "neurodegeneration_downstream_pn.gif"))
FIG_NEURO_TRAJ = animated(
    repo("output", "Neurodegeneration", "images", "neurodegeneration_trajectory.png"),
    repo("output", "Neurodegeneration", "animations", "neurodegeneration_trajectory.gif"))
# No animated counterpart - these stay static.
FIG_NEURO_CMP = repo("output", "Neurodegeneration", "images", "neurodegeneration_comparison.png")
FIG_KEGG_LOCAL = repo("presentation", "NMDAR_Pathway.png")
# Explicit full-resolution perturbation animation (embedded without downscaling).
FIG_PERTURB_ANIM = repo("output", "NMDAR", "animations", "NMDAR_perturbation_animation.gif")


def _downscale_gif(src, dst, max_w):
    """Resize every frame of an animated GIF, preserving timing and looping."""
    from PIL import Image
    im = Image.open(src)
    nw = max_w
    nh = max(1, round(im.height * max_w / im.width))
    try:
        n = im.n_frames
    except Exception:
        n = 1
    frames, durations = [], []
    for i in range(n):
        im.seek(i)
        fr = im.convert("RGB").resize((nw, nh), Image.LANCZOS)
        frames.append(fr.convert("P", palette=Image.ADAPTIVE, colors=256))
        durations.append(im.info.get("duration", 80))
    frames[0].save(dst, save_all=True, append_images=frames[1:],
                   duration=durations, loop=im.info.get("loop", 0),
                   disposal=2, optimize=True)


def prepare_image(path):
    """Return a (possibly downscaled) cached copy of an image for embedding.

    Honours DOWNSCALE_BIG_IMAGES. GIFs are capped at GIF_MAX_WIDTH (all frames
    resized, timing and looping preserved); static images at IMG_MAX_WIDTH.
    Files already within the limit are returned unchanged; originals are never
    overwritten. Cached copies live under presentation/_assets/.
    """
    if not DOWNSCALE_BIG_IMAGES or not os.path.exists(path):
        return path
    from PIL import Image
    ext = os.path.splitext(path)[1].lower()
    is_gif = ext == ".gif"
    max_w = GIF_MAX_WIDTH if is_gif else IMG_MAX_WIDTH
    try:
        with Image.open(path) as probe:
            width = probe.width
    except Exception:
        # File unreadable (e.g. mid-write by the pipeline); embed as-is.
        return path
    if width <= max_w:
        return path
    os.makedirs(ASSETS, exist_ok=True)
    stem = os.path.splitext(os.path.basename(path))[0]
    dst = os.path.join(ASSETS, "%s_w%d%s" % (stem, max_w, ext))
    if os.path.exists(dst) and os.path.getmtime(dst) >= os.path.getmtime(path):
        return dst
    if is_gif:
        _downscale_gif(path, dst, max_w)
    else:
        im = Image.open(path)
        nh = max(1, round(im.height * max_w / im.width))
        if ext in (".jpg", ".jpeg"):
            work = im.convert("RGB")
        elif im.mode in ("RGBA", "LA", "P"):
            work = im.convert("RGBA")
        else:
            work = im.convert("RGB")
        work = work.resize((max_w, nh), Image.LANCZOS)
        if ext in (".jpg", ".jpeg"):
            work.save(dst, quality=88)
        else:
            work.save(dst)
    return dst

# Web image sources kept for reference / fallback captions.
URL_KEGG = "https://www.kegg.jp/pathway/hsa05022"

# Biological-background figures supplied by the user in 'ppt other figures/'.
FIGDIR = repo("presentation", "ppt other figures")


def _as_png(path):
    """python-pptx cannot embed .webp; convert to .png on demand (Pillow)."""
    if path.lower().endswith(".webp"):
        png = os.path.splitext(path)[0] + ".png"
        if not os.path.exists(png) and os.path.exists(path):
            from PIL import Image
            Image.open(path).convert("RGB").save(png)
        return png
    return path


FIG_NMDAR_STRUCT = os.path.join(FIGDIR, "NMDA_receptor.jpg")
FIG_BRAIN_MRI = _as_png(os.path.join(FIGDIR, "fneur-13-834929-g0001.webp"))
FIG_MECHANISM = os.path.join(FIGDIR, "nejmra1708712_f1.jpg")
TITLE_BG = os.path.join(FIGDIR, "slide1_background.jpg")

# --------------------------------------------------------------------------- #
# Palette and fonts
# --------------------------------------------------------------------------- #
NAVY = RGBColor(0x16, 0x33, 0x57)
NAVY_D = RGBColor(0x0E, 0x22, 0x3D)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
GREEN = RGBColor(0x2E, 0x8B, 0x57)   # baseline / healthy
RED = RGBColor(0xC0, 0x39, 0x2B)     # perturbation / disease
INK = RGBColor(0x20, 0x2A, 0x37)
MUTE = RGBColor(0x5B, 0x68, 0x77)
LIGHT = RGBColor(0xEE, 0xF2, 0xF6)
PANEL = RGBColor(0xF4, 0xF7, 0xFA)
PH_FILL = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEARWHITE = RGBColor(0xDC, 0xE6, 0xF0)

BODY = "Calibri"
HEAD = "Calibri"
LIGHTF = "Calibri Light"

SW = Inches(13.333)
SH = Inches(7.5)
MARGIN = Inches(0.62)
CONTENT_W = SW - 2 * MARGIN

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

_counter = {"n": 0}


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #
def new_slide(bg=WHITE):
    slide = prs.slides.add_slide(BLANK)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    _counter["n"] += 1
    return slide


def rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(1.0),
         dash=None, shadow=False, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w
        if dash:
            _set_dash(sp, dash)
    if shadow:
        _soft_shadow(sp)
    return sp


def _set_dash(sp, val="dash"):
    ln = sp.line._get_or_add_ln()
    for el in ln.findall(qn("a:prstDash")):
        ln.remove(el)
    pd = ln.makeelement(qn("a:prstDash"), {"val": val})
    ln.append(pd)


def _soft_shadow(sp):
    spPr = sp._element.spPr
    existing = spPr.find(qn("a:effectLst"))
    if existing is not None:
        spPr.remove(existing)
    xml = (
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:outerShdw blurRad="60000" dist="25000" dir="5400000" rotWithShape="0">'
        '<a:srgbClr val="1A2A3A"><a:alpha val="26000"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
    )
    from pptx.oxml import parse_xml
    spPr.append(parse_xml(xml))


def grad(sp, c1, c2, angle=90):
    """Apply a two-stop linear gradient fill (c1 -> c2) to an existing shape."""
    sp.shadow.inherit = False
    sp.fill.gradient()
    stops = sp.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = c1
    stops[1].position = 1.0
    stops[1].color.rgb = c2
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    sp.line.fill.background()
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb, tf


def add_run(p, text, size, bold=False, color=INK, italic=False, font=BODY):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    return r


def para(tf, first=False, level=0, align=None, before=0, after=6, spacing=1.05):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    if align is not None:
        p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.line_spacing = spacing
    return p


# --------------------------------------------------------------------------- #
# Composite helpers
# --------------------------------------------------------------------------- #
def title_bar(slide, kicker, title):
    grad(rect(slide, 0, 0, SW, Inches(1.18)), NAVY_D, NAVY, angle=90)
    rect(slide, 0, 0, Inches(0.16), Inches(1.18), fill=TEAL)        # left accent tab
    rect(slide, 0, Inches(1.18), SW, Inches(0.055), fill=ACCENT)    # underline
    _, tf = textbox(slide, MARGIN, Inches(0.14), CONTENT_W, Inches(0.94),
                    anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, after=0, spacing=1.0)
    add_run(p, kicker.upper(), 12.5, bold=True, color=NEARWHITE, font=HEAD)
    p2 = para(tf, after=0, before=2, spacing=1.0)
    add_run(p2, title, 27, bold=True, color=WHITE, font=HEAD)


def footer(slide):
    rect(slide, MARGIN, SH - Inches(0.46), CONTENT_W, Pt(0.75), fill=LIGHT)
    _, tf = textbox(slide, MARGIN, SH - Inches(0.44), CONTENT_W, Inches(0.34),
                    anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, after=0, spacing=1.0)
    add_run(p, "Anti-NMDAR Encephalitis  |  Petri Net Modelling", 9, color=MUTE)
    add_run(p, "          Leiden University, LIACS", 9, color=MUTE)
    _, tf2 = textbox(slide, SW - MARGIN - Inches(1.2), SH - Inches(0.44),
                     Inches(1.2), Inches(0.34), anchor=MSO_ANCHOR.MIDDLE)
    pn = para(tf2, first=True, after=0, align=PP_ALIGN.RIGHT, spacing=1.0)
    add_run(pn, str(_counter["n"]), 9, bold=True, color=ACCENT)


def content_slide(kicker, title):
    slide = new_slide(WHITE)
    title_bar(slide, kicker, title)
    footer(slide)
    return slide


def bullet(tf, segments, level=0, first=False, after=7, before=0, size=16.5):
    """segments: list of (text, bold, color). A bullet char is prepended."""
    p = para(tf, first=first, level=level, after=after, before=before, spacing=1.06)
    if level == 0:
        add_run(p, "●  ", size - 4, bold=True, color=ACCENT)
    else:
        add_run(p, "–  ", size, bold=True, color=TEAL)
    for seg in segments:
        text = seg[0]
        bold = seg[1] if len(seg) > 1 else False
        color = seg[2] if len(seg) > 2 else INK
        add_run(p, text, size if level == 0 else size - 1.5, bold=bold, color=color)
    return p


# Set EMBED_REPO_FIGURES = False to revert every repo figure back to a path
# placeholder. Web-sourced images (source_kind="Web image") always stay as URL
# placeholders, since there is no local file to embed.
EMBED_REPO_FIGURES = True


def figure_placeholder(slide, l, t, w, h, source, caption, source_kind="Repo image",
                       downscale=True):
    cap_h = Inches(0.62)
    box_h = h - cap_h
    embed = (EMBED_REPO_FIGURES and source_kind == "Repo image"
             and os.path.exists(source))
    if embed:
        # Embed the real figure, scaled to fit the box while keeping aspect ratio.
        src = prepare_image(source) if downscale else source
        pic = slide.shapes.add_picture(src, l, t)
        nw, nh = pic.width, pic.height
        scale = min(int(w) / nw, int(box_h) / nh)
        pic.width = int(nw * scale)
        pic.height = int(nh * scale)
        pic.left = int(l) + (int(w) - pic.width) // 2
        pic.top = int(t) + (int(box_h) - pic.height) // 2
        pic.line.color.rgb = NEARWHITE
        pic.line.width = Pt(0.75)
        _soft_shadow(pic)
    else:
        box = rect(slide, l, t, w, box_h, fill=PH_FILL, line=ACCENT,
                   line_w=Pt(1.5), dash="dash", shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        box.adjustments[0] = 0.04
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        p = para(tf, first=True, after=4, align=PP_ALIGN.CENTER, spacing=1.0)
        add_run(p, "▦  FIGURE PLACEHOLDER", 12.5, bold=True, color=ACCENT, font=HEAD)
        p2 = para(tf, after=3, align=PP_ALIGN.CENTER, spacing=1.0)
        add_run(p2, "Replace this box with the image below", 10, italic=True, color=MUTE)
        p3 = para(tf, after=0, align=PP_ALIGN.CENTER, spacing=1.02)
        add_run(p3, source_kind + ":  ", 10.5, bold=True, color=INK)
        add_run(p3, source, 10.5, color=ACCENT)
    # caption underneath
    _, ctf = textbox(slide, l, t + box_h + Inches(0.04), w, cap_h,
                     anchor=MSO_ANCHOR.TOP)
    cp = para(ctf, first=True, after=0, align=PP_ALIGN.CENTER, spacing=1.0)
    add_run(cp, caption, 10.5, italic=True, color=MUTE)


SECTION_TOTAL = 5
GHOST = RGBColor(0x20, 0x3D, 0x63)      # subtle watermark number
DOT_OFF = RGBColor(0x33, 0x4E, 0x73)    # inactive progress dot
DIV_FOOT = RGBColor(0x8B, 0x9C, 0xB3)   # muted footer on navy


def section_divider(number, title, subtitle=None):
    slide = new_slide(NAVY)
    grad(rect(slide, 0, 0, SW, SH), NAVY_D, NAVY, angle=45)
    rect(slide, 0, 0, Inches(0.32), SH, fill=ACCENT)
    rect(slide, Inches(0.32), 0, Inches(0.10), SH, fill=TEAL)
    try:
        idx = int(str(number).split()[-1])
    except Exception:
        idx = 0
    # subtle ghost number, lower right
    _, gtf = textbox(slide, Inches(7.4), Inches(2.1), Inches(5.5), Inches(3.4),
                     anchor=MSO_ANCHOR.MIDDLE)
    gp = para(gtf, first=True, after=0, align=PP_ALIGN.RIGHT, spacing=1.0)
    add_run(gp, ("%02d" % idx) if idx else "", 200, bold=True, color=GHOST, font=LIGHTF)
    # heading block
    _, tf = textbox(slide, Inches(1.2), Inches(2.45), Inches(7.6), Inches(2.2),
                    anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, after=4, spacing=1.0)
    add_run(p, number, 30, bold=True, color=TEAL, font=LIGHTF)
    p2 = para(tf, after=0, before=2, spacing=1.0)
    add_run(p2, title, 40, bold=True, color=WHITE, font=HEAD)
    if subtitle:
        p3 = para(tf, before=10, after=0, spacing=1.05)
        add_run(p3, subtitle, 16, color=NEARWHITE)
    # progress dots
    dy = SH - Inches(1.05)
    for i in range(SECTION_TOTAL):
        dx = Inches(1.22) + Inches(0.36) * i
        if i == idx - 1:
            rect(slide, dx, dy, Inches(0.42), Inches(0.14), fill=ACCENT,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        else:
            rect(slide, dx + Inches(0.13), dy, Inches(0.14), Inches(0.14),
                 fill=DOT_OFF, shape=MSO_SHAPE.OVAL)
    # footer
    _, ftf = textbox(slide, Inches(1.2), SH - Inches(0.55), Inches(8.0), Inches(0.34))
    pf = para(ftf, first=True, after=0, spacing=1.0)
    add_run(pf, "Anti-NMDAR Encephalitis  |  Petri Net Modelling", 9, color=DIV_FOOT)
    return slide


def styled_table(slide, l, t, w, headers, rows, col_ratio, font_size=11,
                 head_size=11.5):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h = Inches(0.34)
    gtbl = slide.shapes.add_table(n_rows, n_cols, l, t, w, row_h * n_rows)
    table = gtbl.table
    table.first_row = False
    table.horz_banding = False
    total = sum(col_ratio)
    for j, ratio in enumerate(col_ratio):
        table.columns[j].width = Emu(int(int(w) * ratio / total))
    # header
    for j, htext in enumerate(headers):
        c = table.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        c.margin_top = Pt(2)
        c.margin_bottom = Pt(2)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = c.text_frame
        pp = tf.paragraphs[0]
        pp.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        add_run(pp, htext, head_size, bold=True, color=WHITE, font=HEAD)
    # body
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else PANEL
            c.margin_top = Pt(1)
            c.margin_bottom = Pt(1)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame
            pp = tf.paragraphs[0]
            pp.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            col = INK
            if j == 1:
                col = GREEN
            elif j == 2:
                col = RED
            add_run(pp, str(val), font_size, bold=(j == 0), color=col if j else INK)
    for r in range(n_rows):
        table.rows[r].height = row_h
    return gtbl


def note_panel(slide, l, t, w, h, heading, lines):
    rect(slide, l, t, w, h, fill=PANEL, line=NEARWHITE, line_w=Pt(1.0),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, l, t, Inches(0.10), h, fill=ACCENT)
    _, tf = textbox(slide, l + Inches(0.22), t + Inches(0.12),
                    w - Inches(0.36), h - Inches(0.24))
    p = para(tf, first=True, after=6, spacing=1.0)
    add_run(p, heading, 13.5, bold=True, color=NAVY, font=HEAD)
    for seg in lines:
        bullet(tf, [seg], level=1, after=5, size=13.5)


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def slide_title():
    slide = new_slide(NAVY)
    # Use the base deck's title-slide design as a full-bleed background image
    # (it already carries the title and author text). Falls back to the styled
    # navy layout below if the image is missing.
    if os.path.exists(TITLE_BG):
        pic = slide.shapes.add_picture(prepare_image(TITLE_BG), 0, 0)
        nw, nh = pic.width, pic.height
        scale = max(int(SW) / nw, int(SH) / nh)
        pic.width = int(nw * scale)
        pic.height = int(nh * scale)
        pic.left = (int(SW) - pic.width) // 2
        pic.top = (int(SH) - pic.height) // 2
        return
    rect(slide, 0, 0, Inches(0.32), SH, fill=ACCENT)
    rect(slide, Inches(0.32), 0, Inches(0.10), SH, fill=TEAL)
    _, tf = textbox(slide, Inches(1.1), Inches(1.55), Inches(11.2), Inches(2.7))
    p = para(tf, first=True, after=2, spacing=1.0)
    add_run(p, "BIO-MODELLING PROJECT", 14, bold=True, color=TEAL, font=HEAD)
    p1 = para(tf, before=10, after=2, spacing=1.02)
    add_run(p1, "Modeling Anti-NMDAR Encephalitis", 40, bold=True, color=WHITE, font=HEAD)
    p2 = para(tf, after=0, spacing=1.02)
    add_run(p2, "using Petri Nets", 40, bold=True, color=WHITE, font=HEAD)
    p3 = para(tf, before=14, after=0, spacing=1.08)
    add_run(p3, "A discrete stochastic Petri net of NMDAR signalling and its "
                "downstream neurodegeneration pathway, baseline vs. anti-NMDAR perturbation.",
            16, color=NEARWHITE)
    # author band
    rect(slide, Inches(1.1), Inches(5.55), Inches(10.8), Pt(1.4), fill=TEAL)
    _, atf = textbox(slide, Inches(1.1), Inches(5.7), Inches(11.0), Inches(1.3))
    pa = para(atf, first=True, after=3, spacing=1.0)
    add_run(pa, "Alexa van Thiel (s3617238)   |   Ali Momen (s4993128)   |   "
                "Thijs Vijgeboom (s2648261)", 15, bold=True, color=WHITE)
    pb = para(atf, after=0, spacing=1.0)
    add_run(pb, "LIACS, Leiden University, the Netherlands", 13, color=NEARWHITE)


def slide_outline():
    slide = content_slide("Overview", "Outline")
    items = [
        ("1", "Introduction", "The disease, the gap, the research question"),
        ("2", "Biological Background", "Encephalitis, NMDAR, and the autoimmune mechanism"),
        ("3", "Bio-modeling", "Petri net definitions, decisions, and the two coupled modules"),
        ("4", "Analysis and Results", "Validation of baseline vs. perturbation dynamics"),
        ("5", "Conclusions", "Findings, limitations, and future work"),
    ]
    top = Inches(1.65)
    cx = MARGIN + Inches(0.31)
    # vertical connector behind the circles (timeline look)
    rect(slide, cx - Pt(1.5), Inches(1.96), Pt(3), Inches(4.08),
         fill=RGBColor(0xC9, 0xD6, 0xE2))
    for num, head, sub in items:
        rect(slide, MARGIN, top, Inches(0.62), Inches(0.62), fill=NAVY,
             shape=MSO_SHAPE.OVAL)
        _, ntf = textbox(slide, MARGIN, top, Inches(0.62), Inches(0.62),
                         anchor=MSO_ANCHOR.MIDDLE)
        pn = para(ntf, first=True, after=0, align=PP_ALIGN.CENTER, spacing=1.0)
        add_run(pn, num, 20, bold=True, color=WHITE, font=HEAD)
        _, ttf = textbox(slide, MARGIN + Inches(0.92), top, Inches(10.8),
                         Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
        ph = para(ttf, first=True, after=0, spacing=1.0)
        add_run(ph, head + "   ", 20, bold=True, color=NAVY, font=HEAD)
        add_run(ph, "  " + sub, 13.5, italic=True, color=MUTE)
        top += Inches(1.02)


def slide_intro_disease():
    slide = content_slide("1  Introduction", "An autoimmune attack on the synapse")
    _, tf = textbox(slide, MARGIN, Inches(1.5), Inches(5.55), Inches(5.4))
    bullet(tf, [("Anti-NMDAR encephalitis is an autoimmune disease: IgG "
                 "autoantibodies target the GluN1 subunit of NMDARs.",)], first=True)
    bullet(tf, [("Antibodies crosslink and internalise the receptors, lowering "
                 "surface NMDAR density and reducing Ca2+ entry into the neuron.",)])
    bullet(tf, [("Clinical picture: psychiatric symptoms (psychosis, agitation), "
                 "memory deficits, seizures, and movement disorders.",)])
    bullet(tf, [("The gap: ", True), ("individual pathways are well studied, but no "
                 "integrated model links reduced NMDAR activity to the balance of "
                 "protective vs. damaging pathways over time.",)])
    bullet(tf, [("Petri nets are an established formalism for biological signalling "
                 "and synaptic plasticity, so they fit this question well.",)])
    figure_placeholder(slide, Inches(6.65), Inches(1.5), Inches(6.05), Inches(5.4),
                       FIG_NMDAR_STRUCT,
                       "NMDA receptor at a glutamatergic synapse: GluN1 / GluN2 "
                       "tetramer gating Ca2+ influx. Source: Wikimedia Commons.")


def slide_intro_question():
    slide = content_slide("1  Introduction", "Research question and hypothesis")
    rect(slide, MARGIN, Inches(1.55), CONTENT_W, Inches(1.35), fill=PANEL,
         line=ACCENT, line_w=Pt(1.2), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _, qtf = textbox(slide, MARGIN + Inches(0.25), Inches(1.62),
                     CONTENT_W - Inches(0.5), Inches(1.2), anchor=MSO_ANCHOR.MIDDLE)
    pq = para(qtf, first=True, after=3, spacing=1.0)
    add_run(pq, "Research question", 13, bold=True, color=ACCENT, font=HEAD)
    pq2 = para(qtf, after=0, spacing=1.05)
    add_run(pq2, "How does antibody-mediated NMDAR depletion alter the dynamic "
                 "balance between calcium-dependent survival, oxidative stress, "
                 "and apoptotic signalling?", 17, bold=True, color=NAVY)
    _, tf = textbox(slide, MARGIN, Inches(3.2), CONTENT_W, Inches(3.4))
    bullet(tf, [("Hypothesis:  ", True, TEAL),
                ("reduced Ca2+ influx lowers activation of the CaMKII and NF-kB "
                 "survival / plasticity pathways, while increasing relative flux "
                 "through ROS and CHOP / CASP3, favouring oxidative stress and apoptosis.",)],
           first=True)
    bullet(tf, [("Comparison strategy:  ", True),
                ("a healthy baseline net and an anti-NMDAR perturbation net that "
                 "share identical core conditions, so any divergence comes only from "
                 "the disease-specific machinery.",)])
    bullet(tf, [("Validation plan:  ", True),
                ("simulation of token distributions, baseline vs. perturbed initial "
                 "markings, reachability analysis, and sensitivity analysis "
                 "(supported by Charlie and Snoopy exports).",)])


def slide_bg_encephalitis():
    slide = content_slide("2  Biological Background", "Encephalitis and its autoimmune subtype")
    _, tf = textbox(slide, MARGIN, Inches(1.5), Inches(5.55), Inches(5.4))
    bullet(tf, [("Encephalitis is acute inflammation of the brain parenchyma, "
                 "causing altered consciousness, seizures, and behavioural changes.",)],
           first=True)
    bullet(tf, [("Infectious (IE):  ", True),
                ("pathogen-driven (HSV, enteroviruses, arboviruses).",)])
    bullet(tf, [("Autoimmune (AE):  ", True),
                ("antibodies against neuronal surface or intracellular antigens.",)])
    bullet(tf, [("Often misdiagnosed: the high variance of causative agents leads "
                 "to overlapping, ubiquitous pathology.",)])
    bullet(tf, [("Anti-NMDAR encephalitis is an AE subtype, driven by patient IgG "
                 "against the NMDA receptor - the focus of this project.",)])
    figure_placeholder(slide, Inches(6.65), Inches(1.5), Inches(6.05), Inches(5.4),
                       FIG_BRAIN_MRI,
                       "Brain MRI in anti-NMDAR encephalitis: temporal-lobe and "
                       "white-matter T2 / FLAIR lesions. Source: Front. Neurol. "
                       "2022, Fig. 1 (CC BY).")


def slide_bg_nmdar():
    slide = content_slide("2  Biological Background", "NMDAR and the autoimmune mechanism")
    _, tf = textbox(slide, MARGIN, Inches(1.5), Inches(5.55), Inches(5.4))
    bullet(tf, [("NMDAR:  ", True),
                ("a glutamate-gated ion channel with a voltage-dependent Mg2+ block; "
                 "a heterotetramer of GluN1 + GluN2 subunits.",)], first=True)
    bullet(tf, [("It gates Ca2+ influx that is central to long-term potentiation "
                 "(LTP), synaptic plasticity, and neuronal survival.",)])
    bullet(tf, [("Patients produce IgG autoantibodies mainly against the GluN1 domain.",)])
    bullet(tf, [("Pathogenic cascade:", True, RED)])
    bullet(tf, [("antibody binds extracellular GluN1,",)], level=1)
    bullet(tf, [("surface receptors are crosslinked,",)], level=1)
    bullet(tf, [("the NMDAR-antibody complex is internalised,",)], level=1)
    bullet(tf, [("surface NMDAR density falls -> NMDAR loss of function.",)], level=1)
    figure_placeholder(slide, Inches(6.65), Inches(1.5), Inches(6.05), Inches(5.4),
                       FIG_MECHANISM,
                       "Antibody binding, crosslinking, and internalisation of "
                       "surface NMDARs in anti-NMDAR encephalitis. Source: Dalmau "
                       "& Graus, NEJM 2018, Fig. 1.")


def slide_bg_pathway():
    slide = content_slide("2  Biological Background", "Downstream consequences of NMDAR loss")
    _, tf = textbox(slide, MARGIN, Inches(1.5), Inches(5.3), Inches(5.4))
    bullet(tf, [("Healthy:  ", True, GREEN),
                ("Ca2+ -> CaMKII and CREB / NF-kB -> LTP, gene expression, and "
                 "neuronal survival.",)], first=True)
    bullet(tf, [("NMDAR loss disrupts Ca2+ homeostasis and tips the balance:", True, RED)])
    bullet(tf, [("ER stress -> CHOP (DDIT3), a pro-apoptotic factor,",)], level=1)
    bullet(tf, [("mitochondrial Ca2+ imbalance -> ROS / oxidative stress,",)], level=1)
    bullet(tf, [("convergence on CASP3 -> apoptosis / neurodegeneration,",)], level=1)
    bullet(tf, [("reduced NF-kB signalling -> impaired LTP.",)], level=1)
    bullet(tf, [("Our model is a truncated slice of KEGG hsa05022, keeping only the "
                 "branches downstream of NMDAR dysfunction.",)], before=4)
    figure_placeholder(slide, Inches(6.3), Inches(1.5), Inches(6.4), Inches(5.4),
                       FIG_KEGG_LOCAL,
                       "KEGG neurodegeneration pathway (hsa05022) that our model is a "
                       "part of; web source: " + URL_KEGG,
                       source_kind="Repo image")


def slide_tech_def():
    slide = content_slide("3  Bio-modeling  >  3.1 Technical Definitions",
                          "A discrete stochastic Petri net (SPN)")
    _, tf = textbox(slide, MARGIN, Inches(1.5), Inches(7.05), Inches(5.3))
    bullet(tf, [("We model NMDAR signalling at a glutamatergic synapse as a "
                 "discrete stochastic Petri net.",)], first=True)
    bullet(tf, [("A Petri net is a 5-tuple N = (P, T, F, W, M0):", True)])
    bullet(tf, [("P  places - molecular species or states,",)], level=1)
    bullet(tf, [("T  transitions - biochemical reactions or events,",)], level=1)
    bullet(tf, [("F  flow relation - directed arcs between places and transitions,",)], level=1)
    bullet(tf, [("W  arc weights - stoichiometry,",)], level=1)
    bullet(tf, [("M0  initial marking - tokens per place.",)], level=1)
    bullet(tf, [("Stochastic firing via the Gillespie SSA: each transition fires "
                 "with a mass-action propensity over its input tokens.",)], before=3)
    bullet(tf, [("A transition is enabled when every input place holds at least the "
                 "arc weight; one firing per SSA step; ordinary arcs only.",)])
    note_panel(slide, Inches(8.0), Inches(1.55), Inches(4.7), Inches(4.95),
               "Propensity (mass action)",
               [("a(t) = k . product over input places of  m(p)! / [ (m(p) - w(p))! . w(p)! ]",),
                ("k is the rate constant, m(p) the marking of input place p, w(p) the arc weight.",),
                ("Built in a small, fully introspectable Python Petri-net core.",),
                ("Exported to Snoopy (.spn / .cpn) and Charlie (APNN / ANDL) for external analysis.",)])


def slide_decisions_fig():
    slide = content_slide("3  Bio-modeling  >  3.2 Modeling Decisions",
                          "Two coupled nets sharing cytosolic Ca2+")
    figure_placeholder(slide, MARGIN, Inches(1.5), CONTENT_W, Inches(5.45),
                       FIG_NMDAR_SIDE,
                       "Side-by-side overview: baseline (left) and perturbation "
                       "(right) NMDAR nets with the cumulative Ca2+ curve (animated).",
                       source_kind="Repo image")


def slide_decisions_text():
    slide = content_slide("3  Bio-modeling  >  3.2 Modeling Decisions",
                          "Two coupled nets sharing cytosolic Ca2+")
    _, tf = textbox(slide, MARGIN, Inches(1.55), Inches(6.7), Inches(4.95),
                    anchor=MSO_ANCHOR.MIDDLE)
    bullet(tf, [("The model is split into two coupled SPNs joined by a shared "
                 "cytosolic Ca2+ place - the interface between synaptic signalling "
                 "and downstream neurodegeneration.",)], first=True, size=17, after=12)
    bullet(tf, [("Each module has a baseline and a perturbation version; all are "
                 "simulated for 100 Gillespie SSA runs.",)], size=17, after=12)
    bullet(tf, [("Perturbation adds antibody dynamics (k_on ~ 1e5 /M/s), bivalent "
                 "crosslinking, and lysosomal degradation, progressively losing "
                 "surface NMDARs.",)], size=17, after=12)
    bullet(tf, [("Faithful comparison: identical core initial conditions, so every "
                 "divergence comes from the disease-specific machinery alone.",)], size=17)
    note_panel(slide, Inches(7.65), Inches(1.55), Inches(5.05), Inches(5.0),
               "What the perturbation adds",
               [("NMDAR + Anti-GluN1 IgG -> NMDAR-IgG  (T_Ab_bind)",),
                ("2x NMDAR-IgG -> crosslinked cap  (T_crosslink)",),
                ("NMDAR-EphB2 + IgG -> EphB2 disruption  (Mikasova 2012)",),
                ("crosslinked cap -> endosome  (T_endocytosis)",),
                ("endosome -> lysosome  (receptor pool lost)",),
                ("Net effect: surface NMDAR pool drains -> Ca2+ influx and the "
                 "CaMKII / LTP cascade collapse.",)])


def slide_ca_figs():
    slide = content_slide("3.2.1  Anti-NMDAR Ca2+ Influx Module",
                          "The glutamatergic synapse, healthy and perturbed")
    half = Inches(5.9)
    figure_placeholder(slide, MARGIN, Inches(1.5), half, Inches(5.45),
                       FIG_BASELINE_PN,
                       "Baseline: healthy NMDAR signalling Petri net (animated).",
                       source_kind="Repo image")
    figure_placeholder(slide, MARGIN + half + Inches(0.3), Inches(1.5), half, Inches(5.45),
                       FIG_PERTURB_PN,
                       "Perturbation: anti-NMDAR net with IgG binding, crosslinking, "
                       "endocytosis, and lysosomal degradation (animated).",
                       source_kind="Repo image")


def slide_ca_perturb_full():
    slide = content_slide("3.2.1  Anti-NMDAR Ca2+ Influx Module",
                          "The glutamatergic synapse perturbed")
    figure_placeholder(slide, MARGIN, Inches(1.5), CONTENT_W, Inches(5.45),
                       FIG_PERTURB_ANIM,
                       "Perturbation: anti-NMDAR net with IgG binding, crosslinking, "
                       "endocytosis, and lysosomal degradation (animated, full resolution).",
                       source_kind="Repo image", downscale=False)


def slide_ca_text():
    slide = content_slide("3.2.1  Anti-NMDAR Ca2+ Influx Module",
                          "The glutamatergic synapse, healthy and perturbed")
    _, tf = textbox(slide, MARGIN, Inches(1.45), CONTENT_W, Inches(5.5),
                    anchor=MSO_ANCHOR.MIDDLE)
    bullet(tf, [("Models presynaptic vesicle release, glutamate + glycine binding, "
                 "voltage-dependent Mg2+ block / unblock, channel opening, Ca2+ "
                 "influx and extrusion (PMCA / NCX), and CaMKII / LTP activation.",)],
           first=True, size=19, after=16)
    bullet(tf, [("The perturbation variant adds antibody-mediated internalisation: "
                 "tokens drain from NMDAR (closed) into Degraded (lysosome), lowering "
                 "the Ca2+ influx events that feed the downstream net.",)],
           size=19, after=16)


def slide_ca_table():
    slide = content_slide("3.2.1  Anti-NMDAR Ca2+ Influx Module",
                          "Places and initial markings")
    headers = ["Place", "Baseline", "Perturbation"]
    rows = [
        ["Pre-synaptic vesicle", "8", "8"],
        ["Glutamate (cleft)", "4", "4"],
        ["Glycine / D-Serine", "10", "10"],
        ["Mg2+ (cleft)", "14", "14"],
        ["Ca2+ (extracellular)", "80", "80"],
        ["NMDAR (closed)", "10", "10"],
        ["EphB2 (surface)", "5", "5"],
        ["PSD-95 (free)", "6", "6"],
        ["CaMKII (inactive)", "8", "8"],
        ["Ca2+ (cytosol)", "2", "2"],
        ["Anti-GluN1 IgG", "-", "18"],
        ["Internalised (endosome)", "-", "0"],
        ["Degraded (lysosome)", "-", "0"],
    ]
    styled_table(slide, MARGIN, Inches(1.55), Inches(6.6), headers, rows,
                 col_ratio=[3.0, 1.2, 1.4], font_size=10.5, head_size=11.5)
    note_panel(slide, Inches(7.55), Inches(1.55), Inches(5.15), Inches(4.7),
               "Reading the table",
               [("Both nets share identical core synaptic conditions.",),
                ("The perturbation adds three places (shown as '-' in baseline): "
                 "Anti-GluN1 IgG (18, the patient antibody pool), Internalised "
                 "(endosome), and Degraded (lysosome).",),
                ("Tokens flow NMDAR (closed) -> internalisation -> Degraded, depleting "
                 "the surface pool and the rate of Ca2+ influx.",),
                ("Any downstream divergence is therefore due to antibody-driven "
                 "internalisation, not different starting conditions.",)])


def slide_neuro_figs():
    slide = content_slide("3.2.2  Downstream Apoptotic Module",
                          "From cytosolic Ca2+ to neurodegeneration")
    figure_placeholder(slide, MARGIN, Inches(1.5), CONTENT_W, Inches(5.45),
                       FIG_NEURO_PN,
                       "Downstream Petri nets at t = 0: baseline (left, Ca2+ = 35) "
                       "and perturbation (right, Ca2+ = 15) (animated).",
                       source_kind="Repo image")


def slide_neuro_text():
    slide = content_slide("3.2.2  Downstream Apoptotic Module",
                          "From cytosolic Ca2+ to neurodegeneration")
    _, tf = textbox(slide, MARGIN, Inches(1.45), CONTENT_W, Inches(5.5),
                    anchor=MSO_ANCHOR.MIDDLE)
    bullet(tf, [("This module takes cytosolic Ca2+ as its sole input and models "
                 "three converging pathways:",)], first=True, size=19, after=8)
    bullet(tf, [("Mitochondrial Ca2+ uptake -> loss of membrane potential -> "
                 "cytochrome-c release -> CASP9,",)], level=1, size=17, after=6)
    bullet(tf, [("ER stress -> CHOP (DDIT3),",)], level=1, size=17, after=6)
    bullet(tf, [("ROS generation and scavenging.",)], level=1, size=17, after=12)
    bullet(tf, [("CASP3 is the convergence point: ", True),
                ("it integrates the CASP9, CHOP, and ROS inputs.",)], size=19, after=10)
    bullet(tf, [("Active CASP3 -> apoptosis commitment -> a neurodegeneration firing "
                 "transition that accumulates damage.",)], size=19)


def slide_neuro_table():
    slide = content_slide("3.2.2  Downstream Apoptotic Module",
                          "Perturbed rate constants and markings")
    headers = ["Parameter", "Baseline", "Perturbation"]
    rows = [
        ["Initial Ca2+ (cytosol)", "35", "15"],
        ["ER stress induction (k)", "0.20", "0.50"],
        ["Ca2+-mito uptake (k)", "0.30", "0.15"],
        ["CHOP induction (k)", "0.15", "0.40"],
        ["Apoptosis commitment (k)", "0.22", "0.35"],
        ["Neurodegeneration firing (k)", "0.06", "0.18"],
    ]
    styled_table(slide, MARGIN, Inches(1.7), Inches(6.6), headers, rows,
                 col_ratio=[3.2, 1.2, 1.4], font_size=12.5, head_size=12.5)
    note_panel(slide, Inches(7.55), Inches(1.55), Inches(5.15), Inches(4.7),
               "Why these changes",
               [("Lower initial Ca2+ (15 vs. 35) reflects antibody-mediated "
                 "receptor loss upstream.",),
                ("Higher ER-stress induction (0.50 vs. 0.20) reflects enhanced "
                 "cellular stress sensitivity under NMDAR dysfunction.",),
                ("Reduced Ca2+-mito uptake (0.15 vs. 0.30) follows the lower "
                 "cytosolic Ca2+.",),
                ("All other rate constants are held identical, isolating the effect "
                 "of altered Ca2+ signalling and ER-stress susceptibility.",)])


def slide_val_ca():
    slide = content_slide("4  Analysis and Results  >  4.1 Validation",
                          "Calcium influx: healthy vs. anti-NMDAR")
    _, tf = textbox(slide, MARGIN, Inches(1.5), Inches(4.7), Inches(5.0))
    bullet(tf, [("Healthy signalling (green) sustains higher cytosolic Ca2+ and "
                 "greater cumulative influx than the perturbation (red).",)], first=True)
    bullet(tf, [("After 8 s, cumulative Ca2+ influx in the perturbation falls to ",),
                ("43.1% of baseline", True, RED), (",",)])
    bullet(tf, [("which directly reflects antibody-driven NMDAR dysfunction.",)], level=1)
    bullet(tf, [("This collapse is the input signal handed to the downstream "
                 "apoptotic module.",)], before=4)
    figure_placeholder(slide, Inches(5.55), Inches(1.55), Inches(7.15), Inches(5.05),
                       FIG_CA_INFLUX,
                       "Instantaneous cytosolic Ca2+ (top) and cumulative T_Ca_in "
                       "firings (bottom); perturbation reaches 43.1% of baseline.",
                       source_kind="Repo image")


def slide_val_traj_fig():
    slide = content_slide("4  Analysis and Results  >  4.1 Validation",
                          "Downstream trajectories")
    figure_placeholder(slide, MARGIN, Inches(1.5), CONTENT_W, Inches(5.45),
                       FIG_NEURO_TRAJ,
                       "Five-panel trajectory: CASP3, CHOP, ROS, apoptosis tokens, "
                       "and cumulative neurodegeneration firings (animated).",
                       source_kind="Repo image")


def slide_val_traj_text():
    slide = content_slide("4  Analysis and Results  >  4.1 Validation",
                          "Downstream trajectories")
    _, tf = textbox(slide, MARGIN, Inches(1.45), CONTENT_W, Inches(5.5),
                    anchor=MSO_ANCHOR.MIDDLE)
    bullet(tf, [("A single SSA run over 10 s, baseline (green) vs. perturbation (red).",)],
           first=True, size=19, after=14)
    bullet(tf, [("All markers - active CASP3, CHOP, ROS, apoptosis commitment, and "
                 "cumulative neurodegeneration - are elevated under perturbation.",)],
           size=19, after=14)
    bullet(tf, [("Divergence emerges after t ~ 2 s.",)], size=19, after=14)
    bullet(tf, [("Cumulative neurodegeneration firings: ",),
                ("13 (perturbation) vs. 1 (baseline)", True, RED),
                (" by t = 10 s, a 13-fold increase.",)], size=19)


def slide_val_ensemble():
    slide = content_slide("4  Analysis and Results  >  4.1 Validation",
                          "Ensemble comparison across 100 runs")
    _, tf = textbox(slide, MARGIN, Inches(1.5), Inches(4.7), Inches(5.0))
    bullet(tf, [("Marker accumulation at t = 10 s across 100 SSA runs "
                 "(mean +/- SD).",)], first=True)
    bullet(tf, [("All markers are significantly elevated under perturbation.",)])
    bullet(tf, [("Largest effects:", True)])
    bullet(tf, [("apoptosis commitment, delta = 25.47,",)], level=1)
    bullet(tf, [("neurodegeneration firings, delta = 16.66.",)], level=1)
    bullet(tf, [("This supports the hypothesis: NMDAR depletion shifts the balance "
                 "toward oxidative stress and apoptosis.",)], before=4)
    figure_placeholder(slide, Inches(5.55), Inches(1.55), Inches(7.15), Inches(5.05),
                       FIG_NEURO_CMP,
                       "Mean +/- SD marker counts (left) and the perturbation - "
                       "baseline difference (right) over 100 runs.",
                       source_kind="Repo image")


def slide_conclusions():
    slide = content_slide("5  Conclusions", "What the model shows")
    _, tf = textbox(slide, MARGIN, Inches(1.55), CONTENT_W, Inches(5.0))
    bullet(tf, [("An integrated, two-module stochastic Petri net reproduces the "
                 "anti-NMDAR phenotype:", True)], first=True)
    bullet(tf, [("antibody-driven receptor loss -> reduced Ca2+ influx -> a shift "
                 "from survival (CaMKII / NF-kB / LTP) toward stress and apoptosis "
                 "(ROS, CHOP, CASP3, neurodegeneration).",)], level=1)
    bullet(tf, [("Quantitatively:", True)])
    bullet(tf, [("Ca2+ influx drops to about 43% of baseline,",)], level=1)
    bullet(tf, [("cumulative neurodegeneration rises roughly 13-fold,",)], level=1)
    bullet(tf, [("apoptosis commitment is the largest ensemble effect "
                 "(delta = 25.47).",)], level=1)
    bullet(tf, [("Hypothesis supported: ", True, TEAL),
                ("NMDAR depletion tips the protective / damaging balance toward "
                 "damage, consistent with the anti-NMDAR encephalitis literature.",)],
           before=4)


def slide_limitations():
    slide = content_slide("5  Conclusions", "Limitations and future work")
    note_panel(slide, MARGIN, Inches(1.6), Inches(6.0), Inches(4.9),
               "Limitations",
               [("Mass-action rates are scaled for an 8-10 s simulation window, not "
                 "absolute physiological constants.",),
                ("Sparse quantitative data exist for many transition rates.",),
                ("The net is discrete, stochastic, and uncoloured; the downstream "
                 "perturbation is encoded as a few hand-tuned rate changes.",),
                ("Ordinary arcs only - no inhibitor or read arcs in this version.",)])
    note_panel(slide, Inches(6.85), Inches(1.6), Inches(5.85), Inches(4.9),
               "Future work",
               [("Estimate rates on the critical subgraph and move to a continuous "
                 "(ODE) net where data allow.",),
                ("Add inhibitor / read arcs and coloured tokens for richer logic.",),
                ("Formal analysis in Charlie: P / T-invariants, boundedness, "
                 "reachability, and sensitivity.",),
                ("Titre-dependent antibody dosing to model disease severity.",),
                ("A sensible follow-up: couple the model to treatment scenarios "
                 "(antibody removal / immunotherapy).",)])


def slide_workdist():
    slide = content_slide("Project", "Work distribution between group members")
    note_panel(slide, MARGIN, Inches(1.6), Inches(3.7), Inches(4.9),
               "Group members",
               [("Alexa van Thiel", True),
                ("Ali Momen", True),
                ("Thijs Vijgeboom", True)])
    note_panel(slide, Inches(4.62), Inches(1.6), Inches(8.05), Inches(4.9),
               "Sample task pool - to distribute among members or just write "
               "'equal contribution'",
               [("Biological background - NMDAR function, Ca2+ / LTP coupling, "
                 "oxidative stress, and the rationale for focusing on anti-NMDAR "
                 "encephalitis.",),
                ("Modelling scope and decisions - pathway truncation from KEGG "
                 "hsa05022, the stochastic uncoloured Petri net choice, and the "
                 "baseline vs. perturbation design.",),
                ("Model implementation, simulation, and analysis - the Python "
                 "Petri-net core, Gillespie SSA, figures, and Snoopy / Charlie "
                 "exports.",),
                ("Report and presentation writing was shared across all members.",
                 False, MUTE)])
    _, ntf = textbox(slide, MARGIN, Inches(6.62), CONTENT_W, Inches(0.4))
    pn = para(ntf, first=True, after=0, spacing=1.0)
    add_run(pn, "Note: the tasks above are unassigned - distribute them among the "
                "members for the final version.", 12, italic=True, color=MUTE)


def slide_references():
    slide = content_slide("Appendix", "Key references")
    refs = [
        "Dalmau J., Graus F. (2018). Antibody-Mediated Encephalitis. N. Engl. J. Med. 378(9): 840-851.",
        "Mony L., Paoletti P. (2023). Mechanisms of NMDA Receptor Regulation. Curr. Opin. Neurobiol. 83: 102815.",
        "Hughes E.G. et al. (2010). Cellular and synaptic mechanisms of anti-NMDA receptor encephalitis. J. Neurosci. 30(17): 5866-5875.",
        "Mikasova L. et al. (2012). Disrupted surface cross-talk between NMDA and Ephrin-B2 receptors in anti-NMDA encephalitis. Brain 135(5): 1606-1621.",
        "Moscato E.H. et al. (2014). Acute mechanisms underlying antibody effects in anti-NMDAR encephalitis. Ann. Neurol. 76(1): 108-119.",
        "Hardy S., Robillard P.N. (2005). Petri net modeling and simulation of LTP. Biosystems 82(1): 26-38.",
        "Harding H.P. et al. (2003). An integrated stress response regulates resistance to oxidative stress. Mol. Cell 11(3): 619-633.",
        "Yuan J., Yankner B.A. (2000). Apoptosis in the nervous system. Nature 407: 802-809.",
        "KEGG Pathway Database - Pathways of neurodegeneration (hsa05022). https://www.kegg.jp/pathway/hsa05022",
    ]
    _, tf = textbox(slide, MARGIN, Inches(1.55), CONTENT_W, Inches(5.4))
    first = True
    for r in refs:
        p = para(tf, first=first, after=6, spacing=1.04)
        add_run(p, "–  ", 12.5, bold=True, color=ACCENT)
        add_run(p, r, 12.5, color=INK)
        first = False


# --------------------------------------------------------------------------- #
# Assemble
# --------------------------------------------------------------------------- #
def main():
    slide_title()
    slide_outline()
    section_divider("Section 1", "Introduction",
                    "The disease, the gap it leaves, and the question we ask")
    slide_intro_disease()
    slide_intro_question()
    section_divider("Section 2", "Biological Background",
                    "Encephalitis, the NMDA receptor, and the autoimmune mechanism")
    slide_bg_encephalitis()
    slide_bg_nmdar()
    slide_bg_pathway()
    section_divider("Section 3", "Bio-modeling",
                    "Petri net definitions, modelling decisions, and the two modules")
    slide_tech_def()
    slide_decisions_fig()
    slide_decisions_text()
    slide_ca_figs()
    slide_ca_perturb_full()
    slide_ca_text()
    slide_ca_table()
    slide_neuro_figs()
    slide_neuro_text()
    slide_neuro_table()
    section_divider("Section 4", "Analysis and Results",
                    "Validating baseline against the anti-NMDAR perturbation")
    slide_val_ca()
    slide_val_traj_fig()
    slide_val_traj_text()
    slide_val_ensemble()
    section_divider("Section 5", "Conclusions",
                    "Findings, limitations, and where this could go next")
    slide_conclusions()
    slide_limitations()
    slide_workdist()
    slide_references()

    prs.core_properties.title = "Modeling Anti-NMDAR Encephalitis using Petri Nets"
    prs.core_properties.author = "Alexa van Thiel, Ali Momen, Thijs Vijgeboom"
    prs.save(OUT)
    print("Saved:", OUT)
    print("Slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
